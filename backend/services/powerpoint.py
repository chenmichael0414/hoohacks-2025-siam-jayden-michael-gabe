import os
from google import genai
from dotenv import load_dotenv
from pptx import Presentation
from markdown2 import markdown
from fpdf import FPDF, HTMLMixin
import unicodedata

load_dotenv()
api_key = os.getenv('google_API_key')
client = genai.Client(api_key=api_key)


class MarkdownPDF(FPDF, HTMLMixin):
    pass


def clean_text(text):
    return unicodedata.normalize("NFKD", text).encode("ascii", "ignore").decode("ascii")


def detect_and_write_slides_from_pptx(file_path):
    # 📝 Extract text from PowerPoint
    text = ""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text + " "
        text += "\n\n"

    # 🧠 Prompt for markdown note generation
    prompt = (
        "You will be given raw text extracted from presentation slides. "
        "Your task is to generate well-structured MARKDOWN notes with clean formatting that can be turned into a professional PDF.\n\n"
        "Use Markdown syntax:\n"
        "- `#`, `##`, `###` for headers\n"
        "- `**bold**`, `*italics*` for emphasis\n"
        "- Lists using `-` or `1.`\n"
        "- Code blocks with triple backticks if appropriate\n"
        "Correct any spelling or formatting errors, and ensure the structure is logical and visually organized.\n\n"
        "Text:\n" + text
    )

    # 🔮 Generate notes using Gemini
    response = client.models.generate_content(
        model='gemini-2.0-flash',
        contents=[prompt]
    )
    markdown_notes = response.text

    # 📄 Convert to PDF
    html = markdown(markdown_notes)
    html = clean_text(html)

    pdf = MarkdownPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.write_html(html)
    pdf.output("lecture_notes.pdf")

    return markdown_notes
