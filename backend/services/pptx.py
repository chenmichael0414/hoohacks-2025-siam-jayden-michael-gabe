import os
from google import genai
from dotenv import load_dotenv
from pptx import Presentation
from PIL import Image
import pytesseract
import io

load_dotenv()
api_key = os.getenv('google_API_key')
client = genai.Client(api_key=api_key)


def detect_and_write_slides_from_pptx(file_path):
    text = ""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for i, shape in slide.shapes:
            if hasattr(shape, "image"):
                img_data = slide.shapes[0].image.blob
                img = Image.open(io.BytesIO(img_data))
                text += pytesseract.image_to_string(img) + "\n"
    prompt = (
        "Generate well-formatted and clean notes with titles and headers from the text that you are supplied with."
        "There may be mistakes in formatting and spelling as of right. Please correct them, if necessary."
        "If there is no text provided, do not generate a response. Or instead, generate an empty response.")

    response = client.models.generate_content(
        model='gemini-2.0-flash',
        contents=[prompt, text]
    )
    return response.text
